#!/usr/bin/env python3
"""Extract text/images/notes from a .pptx for HTML slide conversion.

Notes
- Requires: python-pptx
  Install: python3 -m pip install --user python-pptx
- Outputs a JSON manifest plus an assets/ directory with extracted images.

This script is intentionally simple and robust (best-effort extraction).
"""

from __future__ import annotations

import argparse
import json
import os
from pathlib import Path


def extract_pptx(file_path: str, output_dir: str) -> list[dict]:
    from pptx import Presentation  # type: ignore

    prs = Presentation(file_path)
    out_root = Path(output_dir)
    assets_dir = out_root / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    slides_data: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "number": slide_idx,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
        }

        # Notes
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                slide_data["notes"] = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            slide_data["notes"] = ""

        # Shapes
        for shape in slide.shapes:
            # Text
            if getattr(shape, "has_text_frame", False):
                text = (shape.text or "").strip()
                if not text:
                    continue

                # Title heuristic
                if shape == getattr(slide.shapes, "title", None) and not slide_data["title"]:
                    slide_data["title"] = text
                else:
                    slide_data["content"].append({"type": "text", "text": text})

            # Images
            try:
                # 13 = MSO_SHAPE_TYPE.PICTURE
                if getattr(shape, "shape_type", None) == 13 and getattr(shape, "image", None):
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    image_name = f"slide{slide_idx:02d}_img{len(slide_data['images']) + 1:02d}.{image_ext}"
                    image_path = assets_dir / image_name
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)

                    slide_data["images"].append(
                        {
                            "path": f"assets/{image_name}",
                            "width": int(getattr(shape, "width", 0)),
                            "height": int(getattr(shape, "height", 0)),
                        }
                    )
            except Exception:
                # Best-effort: ignore image extraction errors
                pass

        slides_data.append(slide_data)

    return slides_data


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", help="Path to .pptx")
    ap.add_argument("--out", default=".pptx-extract", help="Output directory")
    args = ap.parse_args()

    out_dir = Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)

    slides = extract_pptx(args.pptx, str(out_dir))

    manifest_path = out_dir / "slides.json"
    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(slides, f, indent=2, ensure_ascii=False)

    print(f"Wrote: {manifest_path}")
    print(f"Assets: {out_dir / 'assets'}")


if __name__ == "__main__":
    main()
